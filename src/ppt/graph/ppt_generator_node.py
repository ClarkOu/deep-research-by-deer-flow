# Copyright (c) 2025 Bytedance Ltd. and/or its affiliates
# SPDX-License-Identifier: MIT

import logging
import os
import subprocess
import uuid
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from src.ppt.graph.state import PPTState

logger = logging.getLogger(__name__)


# def ppt_generator_node(state: PPTState):
#     logger.info("Generating ppt file...")
#     # use marp cli to generate ppt file
#     # https://github.com/marp-team/marp-cli?tab=readme-ov-file
#     generated_file_path = os.path.join(
#         os.getcwd(), f"generated_ppt_{uuid.uuid4()}.pptx"
#     )
#     subprocess.run(["marp", state["ppt_file_path"], "-o", generated_file_path])
#     # remove the temp file
#     os.remove(state["ppt_file_path"])
#     logger.info(f"generated_file_path: {generated_file_path}")
#     return {"generated_file_path": generated_file_path}

def extract_markdown_slides(content):
    """将Markdown内容解析为幻灯片结构"""
    # 按幻灯片分隔符(---)分割内容
    slides = re.split(r'\n\s*---\s*\n', content)
    parsed_slides = []
    
    for slide_text_content in slides:
        slide_content = {}
        current_slide_text = slide_text_content
        
        # 提取标题 (# 或 ## 开头的行)
        title_match = re.search(r'^(#+)\s+(.+)$', current_slide_text, re.MULTILINE)
        if title_match:
            level = len(title_match.group(1))
            title = title_match.group(2).strip()
            slide_content['title'] = title
            slide_content['title_level'] = level
            # 从内容中移除标题行
            current_slide_text = re.sub(r'^#+\s+.+$\n?', '', current_slide_text, count=1, flags=re.MULTILINE)
        
        # 新增：提取表格
        table_data = None
        # 匹配完整的表格块：表头行 + 分隔行 + 数据行
        table_pattern = re.compile(r'''
            (                               # 开始捕获完整表格
                ^\|[^\n]*\|[ \t]*\n         # 表头行: | ... |
                ^\|(?:[- :|]*\|)+[ \t]*\n   # 分隔行: | --- | --- | 等
                (?:^\|[^\n]*\|[ \t]*\n)*    # 数据行: | ... | (0行或多行)
            )
        ''', re.MULTILINE | re.VERBOSE)

        table_match = table_pattern.search(current_slide_text)
        if table_match:
            table_str = table_match.group(1).strip()
            table_lines = table_str.split('\n')
            
            # 解析表头
            headers = [h.strip() for h in table_lines[0].strip('|').split('|')]
            
            # 解析数据行
            rows_data = []
            if len(table_lines) > 2:  # 表头、分隔行、数据行
                for line_idx in range(2, len(table_lines)):
                    if table_lines[line_idx].strip():  # 跳过空行
                        row_cells = [cell.strip() for cell in table_lines[line_idx].strip('|').split('|')]
                        if len(row_cells) == len(headers):  # 基本验证
                            rows_data.append(row_cells)
                        else:
                            logger.warning(f"跳过格式错误的表格行: {table_lines[line_idx]}")
            
            if headers:  # 确保找到了表头
                table_data = {'headers': headers, 'rows': rows_data}
                # 从原始内容中移除表格，避免被其他解析逻辑处理
                current_slide_text = current_slide_text.replace(table_match.group(1), '', 1)
        
        slide_content['table'] = table_data
        
        # 新增：提取子标题（### 标题）
        subtitles = []
        subtitle_matches = re.findall(r'^###\s+(.+)$', current_slide_text, re.MULTILINE)
        if subtitle_matches:
            subtitles = [subtitle.strip() for subtitle in subtitle_matches]
            # 从内容中移除子标题行
            current_slide_text = re.sub(r'^###\s+.+$\n?', '', current_slide_text, flags=re.MULTILINE)
        
        slide_content['subtitles'] = subtitles

        # 提取列表项
        bullet_points = []
        list_items = re.findall(r'^[\s]*[-*]\s+(.+)$', current_slide_text, re.MULTILINE)
        if list_items:
            bullet_points = [item.strip() for item in list_items]
        
        # 提取图片链接
        images = []
        img_matches = re.findall(r'!\[([^\]]*)\]\(([^)]+)\)', current_slide_text)
        for alt_text, img_url in img_matches:
            images.append({'alt': alt_text, 'url': img_url})
            
        # 提取段落文本 (不是列表项且不是图片的文本)
        paragraphs = []
        # 移除列表项和图片链接
        cleaned = re.sub(r'^[\s]*[-*]\s+(.+)$', '', current_slide_text, flags=re.MULTILINE)
        cleaned = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', '', cleaned)
        # 分割成段落并过滤空行
        para_lines = [p.strip() for p in cleaned.split('\n') if p.strip()]
        if para_lines:
            paragraphs = para_lines
            
        slide_content['bullets'] = bullet_points
        slide_content['images'] = images
        slide_content['paragraphs'] = paragraphs
        
        parsed_slides.append(slide_content)
        
    return parsed_slides

def apply_slide_design(slide_layout, slide, slide_content):
    """应用幻灯片设计和格式化"""
    # 设置标题
    if 'title' in slide_content:
        title_shape = slide.shapes.title
        title_shape.text = slide_content['title']
        # 格式化标题
        title_text_frame = title_shape.text_frame
        title_text_frame.clear()
        p = title_text_frame.paragraphs[0]
        p.text = slide_content['title']
        p.alignment = PP_ALIGN.CENTER
        
        # 设置标题字体大小和样式
        run = p.runs[0]
        run.font.size = Pt(36) if slide_content.get('title_level', 2) == 1 else Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色
    
    # 添加内容
    content_shape = None
    used_placeholders = []  # 记录已使用的占位符
    
    # 在幻灯片中寻找内容占位符
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 1:  # 内容占位符
            content_shape = shape
            used_placeholders.append(shape)
            break
    
    # 如果找不到，创建一个文本框
    if not content_shape and (slide_content.get('bullets') or slide_content.get('paragraphs')):
        content_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        )
    
    # 处理内容
    if content_shape:
        tf = content_shape.text_frame
        tf.clear()  # 清除默认文本
        tf.word_wrap = True
        
        # 添加项目符号列表
        for point in slide_content.get('bullets', []):
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            
            # 设置正文文本样式
            run = p.runs[0]
            run.font.size = Pt(20)  # 适中的文本大小
            
        # 添加段落文本
        for para in slide_content.get('paragraphs', []):
            if para.strip():  # 确保不是空行
                p = tf.add_paragraph()
                p.text = para
                p.alignment = PP_ALIGN.LEFT
                
                # 设置正文文本样式
                run = p.runs[0]
                run.font.size = Pt(18)  # 稍微小一点的段落文本
    
    # 删除所有未使用的占位符 (新增部分)
    shapes_to_delete = []
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape not in used_placeholders:
            # 跳过标题占位符
            if shape.placeholder_format.type != 0:  # 0是标题占位符
                shapes_to_delete.append(shape)
                
    # 删除未使用的占位符
    for shape in shapes_to_delete:
        try:
            sp = shape._element
            if sp is not None and sp.getparent() is not None:
                sp.getparent().remove(sp)
        except Exception as e:
            logger.warning(f"Failed to remove placeholder: {e}")

def ppt_generator_node(state: PPTState):
    logger.info("Generating editable pptx file...")
    
    # 获取Markdown内容
    with open(state["ppt_file_path"], "r") as f:
        md_content = f.read()
    
    # 解析Markdown成幻灯片结构
    slides_data = extract_markdown_slides(md_content)
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9比例
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 处理每个幻灯片
    for i, slide_content_item in enumerate(slides_data):
        # 空白布局，完全自定义内容
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 创建标题
        if 'title' in slide_content_item:
            is_title_slide = i == 0 or slide_content_item.get('title_level', 2) == 1
            
            top = Inches(0.5)
            left = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_shape.text_frame
            title_frame.word_wrap = True
            
            p = title_frame.add_paragraph()
            p.text = slide_content_item['title']
            p.alignment = PP_ALIGN.CENTER
            
            run = p.runs[0]
            if is_title_slide:
                run.font.size = Pt(44)
            else:
                run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(31, 73, 125)  # 深蓝色
        
        # 定义主要内容区域
        content_area_top = Inches(1.8)
        content_area_left = Inches(0.8)
        content_area_width = Inches(8.4)
        current_y_offset = content_area_top

        # 添加表格 (如果存在)
        if slide_content_item.get('table'):
            table_data = slide_content_item['table']
            headers = table_data['headers']
            rows_data = table_data['rows']
            
            num_table_rows = len(rows_data) + 1  # +1 为表头行
            num_table_cols = len(headers)

            if num_table_rows > 0 and num_table_cols > 0:
                # 估算表格高度
                estimated_table_height = Inches(0.5 * num_table_rows)
                
                table_shape = slide.shapes.add_table(
                    num_table_rows, 
                    num_table_cols, 
                    content_area_left, 
                    content_area_top, 
                    content_area_width, 
                    estimated_table_height
                )
                actual_table = table_shape.table
                
                # 设置列宽 (平均分配)
                col_width_emu = int(content_area_width / num_table_cols)
                for c_idx in range(num_table_cols):
                    actual_table.columns[c_idx].width = col_width_emu

                # 填充表头
                for c_idx, header_text in enumerate(headers):
                    cell = actual_table.cell(0, c_idx)
                    cell.text = header_text
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER  # 表头居中
                    run = p.runs[0]
                    run.font.bold = True
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(0, 0, 0)

                # 填充数据行
                for r_idx, row_cells in enumerate(rows_data):
                    for c_idx, cell_text in enumerate(row_cells):
                        cell = actual_table.cell(r_idx + 1, c_idx)
                        cell.text = cell_text
                        p = cell.text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.LEFT  # 内容左对齐
                        run = p.runs[0]
                        run.font.size = Pt(16)
                        run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 添加其他内容 (项目符号和段落) - 仅当没有表格时
        elif slide_content_item.get('bullets') or slide_content_item.get('paragraphs'):
            # 为不同幻灯片类型设置不同布局
            text_box_top = content_area_top
            text_box_left = content_area_left
            text_box_width = content_area_width
            text_box_height = Inches(3.5)

            if i == 0 and not slide_content_item.get('title_level', 2) == 1:  # 首页特殊布局
                text_box_top = Inches(2.5)
                text_box_left = Inches(1.5)
                text_box_width = Inches(7)
                text_box_height = Inches(2.5)
            
            content_shape = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
            tf = content_shape.text_frame
            tf.word_wrap = True
            
            # 项目符号列表
            for point in slide_content_item.get('bullets', []):
                p = tf.add_paragraph()
                p.text = "• " + point
                p.level = 0
                p.alignment = PP_ALIGN.LEFT
                
                run = p.runs[0]
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(0, 0, 0)
                            
            # 添加段落
            for para in slide_content_item.get('paragraphs', []):
                if para.strip():
                    p = tf.add_paragraph()
                    p.text = para
                    p.alignment = PP_ALIGN.LEFT
                    
                    run = p.runs[0]
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 添加页脚
        if i > 0:  # 不在首页添加
            footer_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.1), Inches(9), Inches(0.3)
            )
            footer_frame = footer_shape.text_frame
            p = footer_frame.add_paragraph()
            p.text = f"第 {i+1} 页"
            p.alignment = PP_ALIGN.RIGHT
            
            run = p.runs[0]
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(128, 128, 128)  # 灰色
    
    # 保存文件
    generated_file_path = os.path.join(
        os.getcwd(), f"generated_ppt_{uuid.uuid4()}.pptx"
    )
    prs.save(generated_file_path)
    
    # 移除临时文件
    os.remove(state["ppt_file_path"])
    
    logger.info(f"generated_file_path: {generated_file_path}")
    return {"generated_file_path": generated_file_path}